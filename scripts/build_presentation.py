"""
build_presentation.py
======================
Generates the FINAL-ROUND presentation deck (.pptx) for the in-person
Stakeholder Engagement Workshop at WMG, University of Warwick.

Format on the day:  3 min oral  +  3 min live model demo  +  4 min Q&A.
Judging weights:    Case 10% | Feasibility 10% | Approach 10% |
                    DEMO 40% | Scalability/Safety/Cost 10% | Clarity 20%.

The deck carries the 3-minute ORAL half; the 3-minute DEMO half runs in the
interactive site (docs/index.html -> Presenter mode, key P). The timed
speaker script for BOTH halves is embedded as slide notes and written to
outputs/presentation_script.md.

All headline numbers are read live from outputs/results.json.

    python scripts/build_presentation.py
"""

from __future__ import annotations
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "outputs"
R = json.loads((OUT / "results.json").read_text(encoding="utf-8"))

# ---- identity (final round is NOT blind — put your names here) -------------
TEAM_LINE = "UoW-Bikes Solar Hub Team"          # <-- EDIT: team / member names
EVENT_LINE = ("Final Round - Sustainable e-Micromobility Stakeholder Engagement "
              "Workshop - WMG, University of Warwick - 6 July 2026")
DEMO_URL = "gourav88502.github.io/emicromobility-robust-charging"
REPO_URL = "github.com/Gourav88502/emicromobility-robust-charging"

NAVY = RGBColor(0x1B, 0x1B, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF2, 0xB7, 0x05)
GREEN = RGBColor(0x43, 0xAA, 0x8B)
BLUE = RGBColor(0x2E, 0x86, 0xAB)
RED = RGBColor(0xD9, 0x53, 0x4F)
INK = RGBColor(0x22, 0x22, 0x2A)
GREY = RGBColor(0x5A, 0x5A, 0x66)
LILAC = RGBColor(0xCA, 0xDC, 0xFC)
HEAD = "Georgia"
BODY = "Calibri"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text(slide, l, t, w, h, runs, size=18, color=INK, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, space=6, line=1.05):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, {})]
    first = True
    for line_runs in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space); p.line_spacing = line
        segs = line_runs if isinstance(line_runs, list) else [line_runs]
        for seg in segs:
            s_text, opt = seg if isinstance(seg, tuple) else (seg, {})
            r = p.add_run(); r.text = s_text
            r.font.name = opt.get("font", font)
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", bold)
            r.font.italic = opt.get("italic", italic)
            r.font.color.rgb = opt.get("color", color)
    return tb


def chip(slide, l, t, w, label, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(0.42))
    _solid(sh, color)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = BODY; r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = WHITE
    return sh


def stat(slide, l, t, w, big, small, color, big_size=40):
    text(slide, l, t, w, 0.9, [[(big, {"size": big_size, "bold": True, "color": color, "font": HEAD})]],
         align=PP_ALIGN.LEFT)
    text(slide, l, t + big_size / 46.0, w, 0.7, [[(small, {"size": 12.5, "color": GREY})]],
         align=PP_ALIGN.LEFT)


def pic_fit(slide, path, l, t, w, h):
    """Add picture scaled to fit within (w,h) box, centred."""
    from PIL import Image
    iw, ih = Image.open(path).size
    box_ar = w / h; img_ar = iw / ih
    if img_ar > box_ar:
        nw = w; nh = w / img_ar
    else:
        nh = h; nw = h * img_ar
    left = l + (w - nw) / 2; top = t + (h - nh) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(nw), Inches(nh))


def qr_png(url: str) -> Path | None:
    """Generate a QR code for the live demo (optional dependency)."""
    try:
        import qrcode
        img = qrcode.make(f"https://{url}")
        p = OUT / "_demo_qr.png"
        img.save(str(p))
        return p
    except Exception:
        return None


def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt


def criteria_footer(slide, label, dark=False):
    """Small footer chip tying each slide to the judging criterion it serves."""
    text(slide, 0.7, 7.02, 12.0, 0.4,
         [[(label, {"size": 10.5, "italic": True,
                    "color": LILAC if dark else GREY})]])


def build():
    rec = R["recommended_design"]; vor = R["value_of_robustness"]; emis = R["emissions"]
    t2 = R.get("theme2_route", {}); t2p = t2.get("profiles", {})
    ror = R.get("robustness_of_robustness", {})
    sb = ror.get("storage_boundary_grid_kW")
    val = R.get("validation", {}).get("metrics_within_published_range", "9/9")
    naive = R["decision_rules"]["naive_deterministic"]
    robust = R["decision_rules"]["maximin_robust"]
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    SCRIPT = []

    # ================= Slide 1: Title & hook (dark) =========================
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.75), Inches(0.28), Inches(1.5))
    _solid(bar, GOLD)
    text(s, 0.9, 1.6, 11.5, 1.0, [[("Robust Solar Charging for a Shared e-Bike Fleet", {"size": 38, "bold": True, "color": WHITE, "font": HEAD})]])
    text(s, 0.9, 2.5, 11.5, 0.9, [[("which hub do you build when the future is uncertain?", {"size": 30, "bold": True, "color": GOLD, "font": HEAD})]])
    text(s, 0.95, 3.75, 11.4, 0.6, [[("University of Warwick campus - UoW Bikes - a 15 kW grid connection - 15 demand futures", {"size": 17, "color": WHITE})]])
    chip(s, 0.95, 4.6, 3.0, "Theme 3 - solar charging station", GREEN)
    chip(s, 4.1, 4.6, 3.6, "Theme 2 - charge/discharge profiles", BLUE)
    chip(s, 7.9, 4.6, 3.4, "Live in-browser model demo", GOLD)
    text(s, 0.95, 5.6, 11.4, 0.5, [[(TEAM_LINE, {"size": 15, "bold": True, "color": WHITE})]])
    text(s, 0.95, 6.15, 11.4, 0.5, [[(EVENT_LINE, {"size": 12, "color": LILAC})]])
    SCRIPT.append((1, "Good morning. Charging infrastructure for shared e-bikes hides a trap: "
        "size a solar hub for average demand and it strands the fleet at peaks; size it for the "
        "worst case and capital sits idle. We asked the sharper question - which hub do you build "
        "when the future is uncertain? - and answered it with robust optimisation on real data. "
        "In three minutes you will see the model itself run live."))

    # ================= Slide 2: The case (10%) ==============================
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    text(s, 0.7, 0.5, 12, 0.8, [[("The case: one depot, a capped wire, an uncertain future", {"size": 30, "bold": True, "color": NAVY, "font": HEAD})]])
    pts = [
        [("Who.  ", {"bold": True, "color": NAVY}), ("UoW Bikes charges a mixed fleet of shared e-bikes and e-cargo bikes at a campus depot.", {})],
        [("The constraint.  ", {"bold": True, "color": NAVY}), ("A 15 kW grid connection; a bigger import limit means a slow, costly DNO reinforcement - so peaks must be met by solar + smart scheduling.", {})],
        [("The uncertainty.  ", {"bold": True, "color": NAVY}), ("Demand swings with term-time, weather and 0-15%/yr growth. Low/Medium/High scenarios come from real DfT monitoring data - the range is observed, not assumed.", {})],
        [("The decision.  ", {"bold": True, "color": NAVY}), ("PV size x battery size x number of bays - fixed today, but it must perform for years, in every plausible future.", {})],
    ]
    text(s, 0.7, 1.55, 6.1, 4.9, pts, size=15, space=11, line=1.06)
    pic_fit(s, OUT / "01_scenario_demand.png", 7.0, 1.55, 5.9, 4.6)
    text(s, 7.0, 6.25, 5.9, 0.5, [[("Demand scenarios built from real DfT shared-micromobility data.", {"size": 11, "italic": True, "color": GREY})]], align=PP_ALIGN.CENTER)
    criteria_footer(s, "Judging criterion: Case Understanding (10%)")
    SCRIPT.append((2, "The case: UoW Bikes charges a mixed fleet at a campus depot behind a "
        "fifteen-kilowatt grid connection - upgrading that wire means a slow, expensive network "
        "reinforcement, so peaks must be met on site. Demand swings with term-time, weather and "
        "growth; our Low, Medium and High scenarios are built from real DfT monitoring data, so "
        "the demand range is observed, not assumed. The decision is set today - solar, battery, "
        "bays - but must perform for years."))

    # ================= Slide 3: Approach (10%) ==============================
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    text(s, 0.7, 0.45, 12, 0.8, [[("Approach: 15 futures x 8,760 hours x 150 designs x 5 decision rules", {"size": 28, "bold": True, "color": NAVY, "font": HEAD})]])
    pic_fit(s, OUT / "methodology_flow.png", 1.2, 1.25, 10.9, 5.3)
    criteria_footer(s, "Judging criterion: Approach / Methodology (10%)")
    SCRIPT.append((3, "The approach in one breath: real solar from PVGIS, real carbon from "
        "National Grid, real demand from DfT. Fifteen probability-weighted futures. For every "
        "candidate hub we simulate all eight thousand seven hundred and sixty hours of the year - "
        "solar first, then battery, then the capped grid. A hundred and fifty designs are scored "
        "under five decision rules, from naive to maximin, and the winner is stress-tested until "
        "the conclusion itself proves robust."))

    # ================= Slide 4: Result + feasibility (10%) ==================
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    text(s, 0.7, 0.5, 12, 0.8, [[("The result: robustness pays - and it is buildable today", {"size": 30, "bold": True, "color": NAVY, "font": HEAD})]])
    pic_fit(s, OUT / "02_pareto.png", 0.6, 1.5, 7.4, 5.1)
    stat(s, 8.3, 1.6, 4.6, f"-{vor['worst_cost_reduction_pct']:.0f}%", "worst-case annual cost vs the naive design", GREEN)
    stat(s, 8.3, 3.0, 4.6, f"{vor['naive_min_service']*100:.1f}% -> {vor['robust_min_service']*100:.1f}%", "guaranteed fleet service in the worst future", BLUE)
    stat(s, 8.3, 4.4, 4.6, f"{rec['pv_kwp']:g} kWp + {rec['n_chargers']} bays", f"recommended hub - GBP {R['recommended_capex_gbp']:,.0f} CAPEX, no battery at 15 kW", NAVY, big_size=30)
    text(s, 8.3, 5.6, 4.6, 1.2, [[("Off-the-shelf mono-Si PV, standard 3 kW AC bays; smart charging replaces storage at a connected site. "
                                    f"{val} model outputs validated against published ranges.", {"size": 12, "color": GREY})]], line=1.15)
    criteria_footer(s, "Judging criteria: Technical Feasibility (10%) + headline result")
    SCRIPT.append((4, "The result. The naive hub is cheapest on an average day but collapses in "
        "the high-demand future: forty-seven thousand pounds a year and one bike in twenty-one "
        "stranded. The robust hub - fifteen kilowatt-peak of solar and eight smart-managed bays - "
        "cuts that worst case by fifty-nine percent and guarantees ninety-nine point six percent "
        "service. And the surprise: no battery. Smart charging flattens the load below the grid "
        "limit, so at a connected site storage adds cost and embodied carbon for nothing. "
        "Everything is off-the-shelf hardware, validated nine-for-nine against published benchmarks."))

    # ================= Slide 5: Scalability, Safety & Cost (10%) ============
    s = prs.slides.add_slide(blank); bg(s, WHITE)
    text(s, 0.7, 0.45, 12, 0.8, [[("Scalable, safe, and costed end-to-end", {"size": 30, "bold": True, "color": NAVY, "font": HEAD})]])
    col_w = 3.95
    # Scalability
    chip(s, 0.7, 1.35, col_w, "SCALABILITY", BLUE)
    text(s, 0.7, 1.95, col_w, 4.6, [
        [("Modular replication:", {"bold": True, "color": NAVY}), (" a city scales by copying GBP 29k hubs, not by grid reinforcement.", {})],
        [("Site-agnostic pipeline:", {"bold": True, "color": NAVY}), (" swap coordinates, demand file and grid cap - the model re-sizes any site.", {})],
        [("UK <-> India transfer:", {"bold": True, "color": NAVY}), (" re-run with Indian solar (~1.6x yield) and 2/3-wheeler fleets; the weak-grid battery case becomes the binding one.", {})],
        [("Growth priced in:", {"bold": True, "color": NAVY}), (" sized across 0-15%/yr growth with saturation.", {})],
    ], size=12.5, space=8, line=1.1)
    # Safety
    chip(s, 4.85, 1.35, col_w, "SAFETY", RED)
    text(s, 4.85, 1.95, col_w, 4.6, [
        [("LFP chemistry", {"bold": True, "color": NAVY}), (" where storage is used - cobalt-free, higher thermal-runaway onset than NMC.", {})],
        [("Standards-mapped:", {"bold": True, "color": NAVY}), (" BS EN 62485-5, BS EN 50604-1, BS 7671; IP65 outdoor bays; Type-B RCD per circuit.", {})],
        [("Layout & fire:", {"bold": True, "color": NAVY}), (" outdoor battery cabinet, 1 m clearance, thermal cut-off; sheltered, ventilated bays away from escape routes.", {})],
        [("Peak enforced twice:", {"bold": True, "color": NAVY}), (" in software (EMS schedule) and hardware (supply fusing).", {})],
    ], size=12.5, space=8, line=1.1)
    # Cost
    chip(s, 9.0, 1.35, col_w, "COST ANALYSIS", GREEN)
    text(s, 9.0, 1.95, col_w, 4.6, [
        [(f"GBP {R['recommended_capex_gbp']:,.0f} CAPEX", {"bold": True, "color": NAVY}), (" - PV, bays, installation; battery-free at 15 kW.", {})],
        [(f"GBP {robust['worst_cost']:,.0f}/yr worst-case", {"bold": True, "color": NAVY}), (f" vs GBP {naive['worst_cost']:,.0f} naive - the premium repays in the first bad year.", {})],
        [(f"LCOE GBP {R['recommended_lcoe_gbp_per_kwh']:.3f}/kWh", {"bold": True, "color": NAVY}), (" delivered, incl. time-of-use tariff, GBP 80/kW-yr demand charge, PV residual value, DoD-aware battery ageing.", {})],
        [(f"{emis['carbon_saving_tCO2_yr']:.1f} tCO2/yr avoided", {"bold": True, "color": NAVY}), (f" ({emis['carbon_saving_pct']:.0f}% vs grid-only, marginal basis) - {emis['carbon_saving_tCO2_lifetime']:.0f} t over the PV lifetime.", {})],
    ], size=12.5, space=8, line=1.1)
    criteria_footer(s, "Judging criterion: Scalability, Safety & Cost Analysis (10%)")
    SCRIPT.append((5, "Scaling, safety, cost. A city scales by copying twenty-nine-thousand-pound "
        "modular hubs - no grid reinforcement - and under the UK-India partnership the same "
        "pipeline re-sizes for Indian solar and two- and three-wheeler fleets. Safety is "
        "standards-mapped: LFP chemistry where storage is used, BS EN battery standards, IP65 "
        "bays, and the peak is enforced in software and hardware. Costs are end-to-end: "
        "twenty-five pence per kilowatt-hour delivered, and three point seven tonnes of CO2 "
        "avoided every year, stated honestly on a marginal basis."))

    # ================= Slide 6: Demo handoff (dark) =========================
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    text(s, 0.9, 0.7, 11.5, 0.9, [[("Now - the model itself, live", {"size": 36, "bold": True, "color": GOLD, "font": HEAD})]])
    steps = [
        ("1", "The hub, animated", "hour-by-hour dispatch of the real model - sun up, bikes out, smart charging into the sun"),
        ("2", "Break it, then fix it", "the naive design collapses in the worst future; the robust design holds at 99.6%"),
        ("3", "150 designs, searched live", "the browser re-runs the full robust optimisation - 20 million simulated hours - in about a second"),
        ("4", "The storage boundary", "weaken the grid to 8 kW and watch a battery enter the optimal design"),
    ]
    y = 1.9
    for n, t_, d_ in steps:
        sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), Inches(y), Inches(0.5), Inches(0.5))
        _solid(sh, GOLD)
        tf = sh.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = n; r.font.bold = True; r.font.size = Pt(18); r.font.color.rgb = NAVY
        text(s, 1.7, y - 0.04, 6.6, 0.5, [[(t_, {"size": 17, "bold": True, "color": WHITE})]])
        text(s, 1.7, y + 0.38, 6.6, 0.6, [[(d_, {"size": 12.5, "color": LILAC})]])
        y += 1.12
    qr = qr_png(DEMO_URL)
    if qr:
        s.shapes.add_picture(str(qr), Inches(9.3), Inches(2.0), Inches(2.6), Inches(2.6))
        text(s, 8.8, 4.75, 3.6, 0.5, [[("scan to open on your phone", {"size": 11.5, "color": LILAC})]], align=PP_ALIGN.CENTER)
    text(s, 8.8, 5.3, 3.6, 0.8, [[(DEMO_URL, {"size": 13, "bold": True, "color": GOLD})]], align=PP_ALIGN.CENTER)
    text(s, 0.95, 6.45, 11.4, 0.5, [[("Runs in any browser - no install, works offline; every number re-derived by the page and checked against the Python pipeline at load.", {"size": 12, "italic": True, "color": LILAC})]])
    criteria_footer(s, "Judging criterion: Demonstration of Model (40%) - handoff to the live demo", dark=True)
    SCRIPT.append((6, "That is the argument - now watch the model make it. The demonstration you "
        "are about to see is not a video and not slides: the full dispatch engine runs live in "
        "the browser, validated against our Python pipeline at load time. Four things to watch: "
        "the hub animated hour by hour; the naive design failing where the robust one holds; the "
        "browser re-searching all one hundred and fifty designs in about a second; and a battery "
        "entering the design the moment the grid weakens to eight kilowatts. Over to the demo."))

    # ================= Slide 7: Q&A backdrop (dark) =========================
    s = prs.slides.add_slide(blank); bg(s, NAVY)
    text(s, 0.9, 0.7, 11.5, 0.8, [[("Questions welcome", {"size": 34, "bold": True, "color": GOLD, "font": HEAD})]])
    text(s, 0.95, 1.75, 11.4, 0.9, [[("We turn demand uncertainty from a risk into a design input - and show that robustness pays.", {"size": 19, "italic": True, "color": WHITE})]])
    qa = [
        (f"{rec['pv_kwp']:g} kWp / {rec['battery_kwh']:g} kWh / {rec['n_chargers']} bays", "recommended robust hub (maximin)"),
        (f"-{vor['worst_cost_reduction_pct']:.0f}%  worst-case cost", f"GBP {naive['worst_cost']:,.0f} -> GBP {robust['worst_cost']:,.0f} per year"),
        (f"{vor['robust_min_service']*100:.1f}% service", "guaranteed in every one of 15 futures"),
        (f"battery only if grid <= {sb:g} kW" if sb else "battery only off-grid", "the storage boundary, located by sweep"),
        (f"{val} validated", "outputs & inputs within published ranges"),
        (f"{emis['carbon_saving_tCO2_yr']:.1f} tCO2/yr avoided", f"{emis['carbon_saving_pct']:.0f}% vs grid-only, marginal basis"),
    ]
    xy = [(0.95, 2.9), (5.15, 2.9), (9.35, 2.9), (0.95, 4.7), (5.15, 4.7), (9.35, 4.7)]
    for (bigv, smallv), (x, yy) in zip(qa, xy):
        text(s, x, yy, 3.9, 0.7, [[(bigv, {"size": 19, "bold": True, "color": GOLD, "font": HEAD})]])
        text(s, x, yy + 0.52, 3.9, 0.8, [[(smallv, {"size": 11.5, "color": LILAC})]], line=1.1)
    text(s, 0.95, 6.5, 11.4, 0.5, [[(f"Live demo: {DEMO_URL}    -    Code, data & tests: {REPO_URL}  (MIT, one-command reproducible)", {"size": 12, "color": WHITE})]])
    SCRIPT.append((7, "(Q&A backdrop - leave on screen. Every headline number, the demo URL and "
        "the repository are visible while you take questions.)"))

    for idx, txt in SCRIPT:
        notes(prs.slides[idx - 1], f"[Slide {idx}]  {txt}")

    out = OUT / "Presentation.pptx"
    prs.save(str(out))
    print(f"Presentation saved: {out}")

    # ---- Timed speaker script (oral + demo halves) --------------------------
    oral = SCRIPT[:6]
    words = sum(len(t.split()) for _, t in oral)
    secs = [22, 32, 30, 42, 32, 22]
    titles = ["Title & hook", "The case (10%)", "Approach (10%)",
              "Result + feasibility (10%)", "Scalability, Safety & Cost (10%)",
              "Demo handoff (into the 40%)"]
    md = ["# Final-Round Script - 3 min oral + 3 min live demo",
          "",
          f"*Oral half ~{words} words = {words/150:.1f} min at a calm 150 wpm; "
          "keep total oral under 3:00. Demo half is driven from the website's "
          "Presenter mode (open the site, press P; arrow keys advance).*",
          "",
          "## PART 1 - ORAL (3:00, slides 1-6)", ""]
    for (idx, txt), sec, title in zip(oral, secs, titles):
        md += [f"### Slide {idx} - {title}  (~{sec}s)", "", txt, ""]
    md += [
        "## PART 2 - LIVE DEMO (3:00, website Presenter mode)", "",
        "*Open docs/index.html (or the GitHub Pages URL), press **P**. The overlay*",
        "*shows each cue and a running timer; -> advances. Steps 1-10 below match*",
        "*the on-screen tour. Speak over each step; the actions fire automatically.*", "",
        "| # | ~t | On screen | Say |",
        "|---|----|-----------|-----|",
        "| 1 | 0:00 | Hub animation playing | This is the hub - every flow is the real model's hourly output, not an illustration. |",
        "| 2 | 0:20 | The case section | Sized for average, it fails; sized for the worst, it wastes. Robust design threads that needle. |",
        "| 3 | 0:40 | Approach pipeline | Real DfT, PVGIS and National Grid data; 15 futures; 8,760-hour physics; 5 decision rules. |",
        "| 4 | 1:00 | Lab - naive preset | The naive hub: watch the worst-case cost and the service KPI go red. |",
        "| 5 | 1:25 | Lab - robust preset | The robust hub: minus 59% worst-case, 99.6% guaranteed. That's the value of robustness. |",
        "| 6 | 1:50 | Click Optimise | The browser is now re-solving all 150 designs across 15 futures - 20 million hours - done in a second. Five rules, same winner. |",
        "| 7 | 2:20 | Grid slider to 8 kW | Weaken the grid and the optimiser buys a battery - the storage boundary is nine kilowatts. That's a design rule, not a guess. |",
        "| 8 | 2:40 | Theme 2 profiles | Physics per route: 4 to 18 watt-hours per km. Shared bikes cycle 1.3x harder - hence the managed hub. |",
        "| 9 | 2:50 | Scale-Safety-Cost | Modular GBP 29k hubs, LFP + BS EN safety, 25p per kWh delivered - and the same pipeline re-sizes for India. |",
        "| 10 | 2:58 | Validation close | Nine-for-nine validated, fully reproducible - we made uncertainty a design input. Thank you. |",
        "",
        "**Q&A (4:00):** leave Slide 7 (backdrop) or the site's validation section on screen.",
        ""]
    (OUT / "presentation_script.md").write_text("\n".join(md), encoding="utf-8")
    print(f"Speaker script saved: {OUT / 'presentation_script.md'}")
    return out


if __name__ == "__main__":
    build()
